#!/usr/bin/env python3
"""报表生成脚本 — 表格、图表、透视表、PPT"""

import argparse
import sys
from pathlib import Path


def table_from_csv(input_file, output=None):
    """从 CSV 生成格式化 Excel 表格"""
    import csv
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill

    wb = Workbook()
    ws = wb.active
    ws.title = "报表"

    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True)

    with open(input_file, encoding="utf-8") as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader):
            for j, val in enumerate(row):
                cell = ws.cell(row=i + 1, column=j + 1, value=val)
                if i == 0:
                    cell.fill = header_fill
                    cell.font = header_font

    output = output or str(Path(input_file).with_suffix(".xlsx"))
    wb.save(output)
    print(f"已生成表格: {output}")


def chart_from_csv(input_file, chart_type="bar", title=None):
    """从 CSV 生成图表"""
    import csv
    from openpyxl import Workbook, load_workbook
    from openpyxl.chart import BarChart, LineChart, PieChart
    from openpyxl.chart.series import DataPoint
    from openpyxl.chart.label import DataLabelList

    wb = Workbook()
    ws = wb.active

    with open(input_file, encoding="utf-8") as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader):
            for j, val in enumerate(row):
                try:
                    ws.cell(row=i + 1, column=j + 1, value=float(val))
                except (ValueError, TypeError):
                    ws.cell(row=i + 1, column=j + 1, value=val)

    chart_map = {"bar": BarChart, "line": LineChart, "pie": PieChart}
    chart_cls = chart_map.get(chart_type, BarChart)
    chart = chart_cls()
    chart.title = title or "数据图表"

    ws.add_chart(chart, "E2")
    print("图表已添加到表格")

    output = str(Path(input_file).with_suffix(".xlsx"))
    wb.save(output)
    print(f"已保存: {output}")


def pivot(input_file, rows, cols, values):
    """生成数据透视表"""
    import csv
    import pandas as pd

    df = pd.read_csv(input_file)
    pivot = pd.pivot_table(df, values=values, index=rows, columns=cols, aggfunc="sum", fill_value=0)

    print("\n=== 数据透视表 ===")
    print(pivot.to_string())

    output = str(Path(input_file).with_suffix(".xlsx"))
    pivot.to_excel(output)
    print(f"\n已导出: {output}")


def summary(input_file):
    """数据统计摘要"""
    import csv
    import pandas as pd

    df = pd.read_csv(input_file)

    print(f"📊 数据概览: {Path(input_file).name}")
    print(f"  行数: {len(df)}, 列数: {len(df.columns)}")
    print(f"  列名: {', '.join(df.columns)}")
    print("\n数值列统计:")
    numeric_cols = df.select_dtypes(include="number").columns
    if len(numeric_cols) > 0:
        print(df[numeric_cols].describe().to_string())
    else:
        print("  无数值列")


def ppt_from_data(input_file, template=None, output=None):
    """从数据生成 PPT"""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    if template and Path(template).exists():
        prs = Presentation(template)
    else:
        prs = Presentation()

    import csv
    import pandas as pd
    df = pd.read_csv(input_file)

    # 标题页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "数据报表"
    slide.placeholders[1].text = f"基于 {Path(input_file).name}"

    # 数据页
    for chunk_start in range(0, len(df), 10):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
        table_shape = slide.shapes.add_table(
            min(11, len(df) - chunk_start + 1),
            len(df.columns),
            Inches(0.5), Inches(1), Inches(9), Inches(5)
        )
        table = table_shape.table

        headers = list(df.columns)
        for j, h in enumerate(headers):
            table.cell(0, j).text = str(h)

        for i, row in df.iloc[chunk_start:chunk_start + 10].iterrows():
            for j, val in enumerate(row):
                table.cell(i - chunk_start + 1, j).text = str(val)

    output = output or "report.pptx"
    prs.save(output)
    print(f"已生成 PPT: {output}")


def main():
    parser = argparse.ArgumentParser(description="报表生成工具")
    subparsers = parser.add_subparsers(dest="operation", required=True)

    # table
    p_table = subparsers.add_parser("table")
    p_table.add_argument("--input", required=True)
    p_table.add_argument("--output")

    # chart
    p_chart = subparsers.add_parser("chart")
    p_chart.add_argument("--input", required=True)
    p_chart.add_argument("--type", default="bar", choices=["bar", "line", "pie"])
    p_chart.add_argument("--title")

    # pivot
    p_pivot = subparsers.add_parser("pivot")
    p_pivot.add_argument("--input", required=True)
    p_pivot.add_argument("--rows", required=True)
    p_pivot.add_argument("--cols", required=True)
    p_pivot.add_argument("--values", required=True)

    # summary
    p_summary = subparsers.add_parser("summary")
    p_summary.add_argument("--input", required=True)

    # ppt
    p_ppt = subparsers.add_parser("ppt")
    p_ppt.add_argument("--input", required=True)
    p_ppt.add_argument("--template")
    p_ppt.add_argument("--output")

    args = parser.parse_args()

    if args.operation == "table":
        table_from_csv(args.input, args.output)
    elif args.operation == "chart":
        chart_from_csv(args.input, args.type, args.title)
    elif args.operation == "pivot":
        pivot(args.input, args.rows, args.cols, args.values)
    elif args.operation == "summary":
        summary(args.input)
    elif args.operation == "ppt":
        ppt_from_data(args.input, args.template, args.output)


if __name__ == "__main__":
    main()
