#!/usr/bin/env python3
"""大纲转 PPT — 从 Markdown/Word 大纲生成演示文稿"""

import argparse
import re
import sys
from pathlib import Path


def parse_outline(filepath):
    """解析大纲，返回层级结构列表 [(level, text), ...]"""
    src = Path(filepath)
    outline = []

    if src.suffix == ".md":
        for line in src.read_text(encoding="utf-8").split("\n"):
            stripped = line.strip()
            if not stripped:
                continue
            m = re.match(r"^(#{1,6})\s+(.*)", stripped)
            if m:
                level = len(m.group(1))
                outline.append((level, m.group(2).strip()))

    elif src.suffix == ".docx":
        from docx import Document
        doc = Document(filepath)
        for para in doc.paragraphs:
            if para.style.name.startswith("Heading"):
                level = int(para.style.name.replace("Heading ", ""))
                outline.append((level, para.text.strip()))
    else:
        print(f"不支持的文件格式: {src.suffix}", file=sys.stderr)
        sys.exit(1)

    return outline


def build_ppt_from_outline(outline, output, template=None):
    """根据大纲生成 PPT"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    if template and Path(template).exists():
        prs = Presentation(template)
    else:
        prs = Presentation()

    # 策略：最小的 2 个层级作为 "分页级"（如 # 和 ##），
    # 分页级标题 = 新幻灯片，其余 = 内容项
    if not outline:
        print("未找到大纲内容", file=sys.stderr)
        return

    levels = sorted(set(l for l, _ in outline))
    # 如果有多个层级，用前两个作为"分页级"；否则全部归为一页
    split_levels = set(levels[:2]) if len(levels) >= 2 else set(levels)

    groups = []
    current_group = None

    for level, text in outline:
        if level in split_levels:
            if current_group:
                groups.append(current_group)
            current_group = {"title": text, "items": []}
        elif current_group is not None:
            current_group["items"].append((level, text))
    if current_group:
        groups.append(current_group)

    # 生成幻灯片
    for i, group in enumerate(groups):
        # 内容幻灯片
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        slide.shapes.title.text = group["title"]

        tf = slide.placeholders[1].text_frame
        tf.clear()

        if group["items"]:
            for j, (level, text) in enumerate(group["items"]):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = text
                p.level = min(level - 2, 3)
        else:
            tf.paragraphs[0].text = "（无详细内容）"

    # 如果没有内容，至少加一页标题
    if not groups and outline:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = outline[0][1]

    prs.save(output)
    print(f"已生成 PPT: {output} ({len(groups)} 页内容幻灯片)")


def main():
    parser = argparse.ArgumentParser(description="大纲转 PPT")
    parser.add_argument("file", help="大纲文件 (.md 或 .docx)")
    parser.add_argument("--output", default=None, help="输出 PPT 路径")
    parser.add_argument("--template", help="PPT 模板路径 (.pptx)")
    args = parser.parse_args()

    output = args.output or str(Path(args.file).with_suffix(".pptx"))
    outline = parse_outline(args.file)

    if not outline:
        print("未找到大纲标题（需要 # Heading 或 Word 的标题样式）", file=sys.stderr)
        sys.exit(1)

    print(f"找到 {len(outline)} 个标题节点:")
    for level, text in outline:
        print(f"  {'  ' * (level - 1)}# {text}")

    build_ppt_from_outline(outline, output, args.template)


if __name__ == "__main__":
    main()
